"""
===========================================================
PowerPoint Processing Demo
Libraries
1. python-pptx
2. Aspose.Slides
3. pptx2md
===========================================================
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ------------------------------------------------------------
# Change this if your filename is different
# ------------------------------------------------------------

FILE = "Copy of Stocks Trading Business Plan.pptx"

# ============================================================
# 1. PYTHON-PPTX
# ============================================================

print("=" * 70)
print("1. PYTHON-PPTX")
print("=" * 70)

prs = Presentation(FILE)

print("\nPresentation Information")
print("---------------------------")

print("Number of Slides :", len(prs.slides))

print("Slide Width :", prs.slide_width)

print("Slide Height :", prs.slide_height)

# ------------------------------------------------------------
# Read Every Slide
# ------------------------------------------------------------

for slide_no, slide in enumerate(prs.slides, start=1):

    print("\n")
    print("=" * 60)
    print(f"Slide {slide_no}")
    print("=" * 60)

    image_count = 0
    table_count = 0
    chart_count = 0

    for shape in slide.shapes:

        print("\nShape Type :", shape.shape_type)

        # ---------------- TEXT ----------------

        if hasattr(shape, "text"):

            text = shape.text.strip()

            if text:

                print("Text :")

                print(text)

        # ---------------- IMAGE ----------------

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

            image_count += 1

            image = shape.image

            ext = image.ext

            filename = f"slide{slide_no}_image{image_count}.{ext}"

            with open(filename, "wb") as f:

                f.write(image.blob)

            print("Image Saved :", filename)

        # ---------------- TABLE ----------------

        if shape.has_table:

            table_count += 1

            print("\nTable Found")

            table = shape.table

            for row in table.rows:

                print([cell.text for cell in row.cells])

        # ---------------- CHART ----------------

        if shape.has_chart:

            chart_count += 1

            print("\nChart Found")

            chart = shape.chart

            print("Chart Type :", chart.chart_type)

    print("\nSummary")

    print("Images :", image_count)

    print("Tables :", table_count)

    print("Charts :", chart_count)


print("\nDone.")